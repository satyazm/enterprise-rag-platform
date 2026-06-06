from pptx import Presentation


def load_pptx(path: str) -> list[dict]:
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append({
                "content": "\n".join(texts),
                "metadata": {"slide": i + 1, "source_type": "pptx"},
            })
    return slides
